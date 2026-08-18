"""
SIH 2026 Presentation Generator: AI-Based Smart Fertilizer Recommendation System
================================================================================
Generates both:
  1. SIH_2026_Smart_Fertilizer_Recommendation.pptx (Editable 16:9 Widescreen PPTX)
  2. SIH_2026_Smart_Fertilizer_Recommendation.pdf (Export-ready 6-slide submission PDF)

Conforms 100% to the official Smart India Hackathon 2026 Sprint Template.
"""

import os
from PIL import Image, ImageDraw, ImageFont
import reportlab
from reportlab.lib.pagesizes import landscape, letter
from reportlab.pdfgen import canvas

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def create_pptx_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    logo_path = 'sih_presentation/extracted_assets/page_1_img_1_Image26.png'
    brain_logo_path = 'sih_presentation/extracted_assets/page_1_img_0_Image7.png'

    # -------------------------------------------------------------
    # SLIDE 1: TITLE PAGE
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)

    # Top Right Logo
    if os.path.exists(logo_path):
        slide1.shapes.add_picture(logo_path, Inches(11.8), Inches(0.4), width=Inches(1.1))

    # Header: SMART INDIA HACKATHON 2026
    tx1 = slide1.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(10.5), Inches(0.8))
    tf1 = tx1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "SMART INDIA HACKATHON 2026"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.bold = True
    p1.font.size = Pt(36)
    p1.font.color.rgb = RGBColor(15, 41, 66)

    # Subheader: TITLE PAGE / Project Title
    tx_sub = slide1.shapes.add_textbox(Inches(1.0), Inches(1.4), Inches(10.5), Inches(0.6))
    tf_sub = tx_sub.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "TITLE PAGE"
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.font.bold = True
    p_sub.font.size = Pt(26)
    p_sub.font.color.rgb = RGBColor(30, 58, 138)

    # Project Title Card (Green / Slate Accent)
    title_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(7.4), Inches(1.3))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = RGBColor(240, 253, 244)
    title_box.line.color.rgb = RGBColor(34, 197, 94)
    title_box.line.width = Pt(1.5)
    tf_tbox = title_box.text_frame
    tf_tbox.word_wrap = True
    p_t = tf_tbox.paragraphs[0]
    p_t.text = "AI-Based Smart Fertilizer Recommendation System"
    p_t.font.bold = True
    p_t.font.size = Pt(22)
    p_t.font.color.rgb = RGBColor(21, 128, 61)
    p_t2 = tf_tbox.add_paragraph()
    p_t2.text = "Precision Nutrient Optimization & Decision-Support Engine for Sustainable Agriculture"
    p_t2.font.size = Pt(13)
    p_t2.font.color.rgb = RGBColor(71, 85, 105)

    # Left Info Bullets (Exact SIH format)
    info_box = slide1.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(7.4), Inches(2.2))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True

    p_ps = tf_info.paragraphs[0]
    p_ps.text = "• Problem Statement ID – [INSERT ACTUAL PS ID]"
    p_ps.font.bold = True
    p_ps.font.size = Pt(18)
    p_ps.font.color.rgb = RGBColor(15, 23, 42)

    p_th = tf_info.add_paragraph()
    p_th.text = "• Theme – Agriculture, FoodTech & Rural Development"
    p_th.font.bold = True
    p_th.font.size = Pt(18)
    p_th.font.color.rgb = RGBColor(15, 23, 42)

    p_tm = tf_info.add_paragraph()
    p_tm.text = "• Team Name – [INSERT REGISTERED TEAM NAME]"
    p_tm.font.bold = True
    p_tm.font.size = Pt(18)
    p_tm.font.color.rgb = RGBColor(15, 23, 42)

    # Right Brain Graphic
    if os.path.exists(brain_logo_path):
        slide1.shapes.add_picture(brain_logo_path, Inches(8.8), Inches(2.2), width=Inches(3.8))

    # Bottom Visual Flow
    flow_box = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.3), Inches(11.73), Inches(0.65))
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = RGBColor(241, 245, 249)
    flow_box.line.color.rgb = RGBColor(203, 213, 225)
    tf_flow = flow_box.text_frame
    p_fl = tf_flow.paragraphs[0]
    p_fl.text = "🌱 SOIL CHEMISTRY (N,P,K,pH,Micros)  ➔  🤖 ML SOFT-VOTING ENSEMBLE  ➔  🧪 PRECISE FERTILIZER DOSAGE  ➔  🌾 SUSTAINABLE CROP YIELD"
    p_fl.alignment = PP_ALIGN.CENTER
    p_fl.font.bold = True
    p_fl.font.size = Pt(12)
    p_fl.font.color.rgb = RGBColor(15, 23, 42)

    # Helper function for common header in slides 2-6
    def setup_slide_header(slide, title_text, template_num):
        # Top-Left Team Name Badge (Oval)
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.3), Inches(1.6), Inches(0.85))
        oval.fill.background()
        oval.line.color.rgb = RGBColor(0, 0, 0)
        oval.line.width = Pt(1.5)
        tf_ov = oval.text_frame
        p_ov = tf_ov.paragraphs[0]
        p_ov.text = "Your\nTeam\nName"
        p_ov.alignment = PP_ALIGN.CENTER
        p_ov.font.size = Pt(11)
        p_ov.font.bold = True
        p_ov.font.color.rgb = RGBColor(0, 0, 0)

        # Center Title
        tx_title = slide.shapes.add_textbox(Inches(2.5), Inches(0.4), Inches(8.3), Inches(0.8))
        tf_t = tx_title.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.alignment = PP_ALIGN.CENTER
        p_t.font.bold = True
        p_t.font.size = Pt(28)
        p_t.font.color.rgb = RGBColor(0, 0, 0)

        # Top-Right Logo
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(11.8), Inches(0.3), width=Inches(1.1))

        # Bottom-Left Footnote
        tx_fn = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(5.0), Inches(0.4))
        tf_fn = tx_fn.text_frame
        p_fn = tf_fn.paragraphs[0]
        p_fn.text = f"@SIH Idea submission- Template {template_num}"
        p_fn.font.size = Pt(10)
        p_fn.font.color.rgb = RGBColor(100, 116, 139)

    # -------------------------------------------------------------
    # SLIDE 2: IDEA TITLE / PROPOSED SOLUTION
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    setup_slide_header(slide2, "IDEA TITLE", 2)

    # Subtitle Header
    tx_sub2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5))
    tf_sub2 = tx_sub2.text_frame
    p_sub2 = tf_sub2.paragraphs[0]
    p_sub2.text = "❖ Proposed Solution (Describe your Idea/Solution/Prototype)"
    p_sub2.font.bold = True
    p_sub2.font.size = Pt(18)
    p_sub2.font.color.rgb = RGBColor(30, 58, 138)

    # 3-Column Cards Layout: PROBLEM | PROPOSED SOLUTION | INNOVATION
    # Card 1: Problem (Red border)
    c1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.9), Inches(3.9), Inches(3.9))
    c1.fill.solid()
    c1.fill.fore_color.rgb = RGBColor(254, 242, 242)
    c1.line.color.rgb = RGBColor(239, 68, 68)
    tf_c1 = c1.text_frame
    tf_c1.word_wrap = True
    p_c1_h = tf_c1.paragraphs[0]
    p_c1_h.text = "⚠️ THE PROBLEM"
    p_c1_h.font.bold = True
    p_c1_h.font.size = Pt(14)
    p_c1_h.font.color.rgb = RGBColor(185, 28, 28)

    bullets1 = [
        "Blind fertilizer application without actual soil nutrient knowledge.",
        "Unbalanced NPK dosage increases input cost and lowers ROI.",
        "Excessive chemicals cause soil acidification & water pollution.",
        "Manual interpretation of 12-parameter soil health cards is complex."
    ]
    for b in bullets1:
        p = tf_c1.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(51, 65, 85)

    # Card 2: Proposed Solution (Blue border)
    c2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.7), Inches(1.9), Inches(3.9), Inches(3.9))
    c2.fill.solid()
    c2.fill.fore_color.rgb = RGBColor(239, 246, 255)
    c2.line.color.rgb = RGBColor(59, 130, 246)
    tf_c2 = c2.text_frame
    tf_c2.word_wrap = True
    p_c2_h = tf_c2.paragraphs[0]
    p_c2_h.text = "💡 PROPOSED SOLUTION"
    p_c2_h.font.bold = True
    p_c2_h.font.size = Pt(14)
    p_c2_h.font.color.rgb = RGBColor(29, 78, 216)

    bullets2 = [
        "AI/ML decision-support system analyzing 12+ soil parameters.",
        "Crop-aware intelligence matching crop-specific stoichiometric demand.",
        "Recommends exact fertilizer product + kg/ha application dosage.",
        "Generates transparent, explainable advice for Indian farmers."
    ]
    for b in bullets2:
        p = tf_c2.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(51, 65, 85)

    # Card 3: Innovation (Green border)
    c3 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(1.9), Inches(3.9), Inches(3.9))
    c3.fill.solid()
    c3.fill.fore_color.rgb = RGBColor(240, 253, 244)
    c3.line.color.rgb = RGBColor(34, 197, 94)
    tf_c3 = c3.text_frame
    tf_c3.word_wrap = True
    p_c3_h = tf_c3.paragraphs[0]
    p_c3_h.text = "⭐ INNOVATION & UNIQUENESS"
    p_c3_h.font.bold = True
    p_c3_h.font.size = Pt(14)
    p_c3_h.font.color.rgb = RGBColor(21, 128, 61)

    bullets3 = [
        "Integrates Liebig's Law of the Minimum across 6 essential micronutrients.",
        "Soft-Voting Ensemble combining Trees (RF, ET), Boosting & Deep MLP.",
        "Real soil benchmark validation across 10.85M national records.",
        "Zero-jargon, actionable outputs with weather spray-safety indexing."
    ]
    for b in bullets3:
        p = tf_c3.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(51, 65, 85)

    # Bottom Pipeline Strip
    pipe_box = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.85))
    pipe_box.fill.solid()
    pipe_box.fill.fore_color.rgb = RGBColor(248, 250, 252)
    pipe_box.line.color.rgb = RGBColor(203, 213, 225)
    tf_pipe = pipe_box.text_frame
    p_pipe = tf_pipe.paragraphs[0]
    p_pipe.text = "🔬 SYSTEM PIPELINE FLOW:\n[Soil Test Data + Crop] ➔ [Data Preprocessing & Scaling] ➔ [Trained ML Meta-Ensemble] ➔ [Fertilizer Type + Dosage kg/ha] ➔ [Farmer-Friendly Action Plan]"
    p_pipe.alignment = PP_ALIGN.CENTER
    p_pipe.font.bold = True
    p_pipe.font.size = Pt(11)
    p_pipe.font.color.rgb = RGBColor(15, 23, 42)

    # -------------------------------------------------------------
    # SLIDE 3: TECHNICAL APPROACH
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    setup_slide_header(slide3, "TECHNICAL APPROACH", 3)

    # Left Box: System Architecture Flowchart
    arch_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(6.8), Inches(5.4))
    arch_box.fill.solid()
    arch_box.fill.fore_color.rgb = RGBColor(248, 250, 252)
    arch_box.line.color.rgb = RGBColor(59, 130, 246)
    tf_arch = arch_box.text_frame
    tf_arch.word_wrap = True
    p_a_h = tf_arch.paragraphs[0]
    p_a_h.text = "🏗️ END-TO-END SYSTEM ARCHITECTURE"
    p_a_h.font.bold = True
    p_a_h.font.size = Pt(14)
    p_a_h.font.color.rgb = RGBColor(30, 58, 138)

    flow_steps = [
        "1. FARMER / SOIL INPUT: Soil N, P, K, pH, OC, EC, Zn, B, S, Fe, Mn, Cu + Selected Crop",
        "2. WEB INTERFACE: Responsive Django web portal with automated input validation",
        "3. PREPROCESSING PIPELINE: Outlier sanitation & RobustScaler (zero data leakage)",
        "4. FEATURE ENGINEERING: 44 stoichiometric ratios, Liebig quotients & buffer capacity",
        "5. ML META-ENSEMBLE: Soft-voting meta-classifier (Random Forest, Extra Trees, HistGradientBoosting, MLP)",
        "6. STOICHIOMETRIC CALCULATOR: Exact physical kg/ha nutrient balance calculation",
        "7. EXPLAINABLE DASHBOARD: Clear justification report + weather spray advisory"
    ]
    for step in flow_steps:
        p = tf_arch.add_paragraph()
        p.text = step
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(30, 41, 59)

    # Right Top Box: Technology Stack
    tech_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.4), Inches(5.3), Inches(2.9))
    tech_box.fill.solid()
    tech_box.fill.fore_color.rgb = RGBColor(239, 246, 255)
    tech_box.line.color.rgb = RGBColor(37, 99, 235)
    tf_tech = tech_box.text_frame
    tf_tech.word_wrap = True
    p_t_h = tf_tech.paragraphs[0]
    p_t_h.text = "⚙️ TECHNOLOGY STACK"
    p_t_h.font.bold = True
    p_t_h.font.size = Pt(14)
    p_t_h.font.color.rgb = RGBColor(29, 78, 216)

    techs = [
        "• Frontend: HTML5, Modern Responsive CSS, JavaScript (Async API fetch)",
        "• Backend: Django Framework, Python 3.14 REST API",
        "• Machine Learning: Scikit-learn, NumPy, Pandas, Joblib",
        "• ML Ensemble: Random Forest (250) + Extra Trees (250) + HistGradientBoosting (250) + MLP",
        "• Database: SQLite (Indexed National Soil Health DB + Farm Plots)"
    ]
    for t in techs:
        p = tf_tech.add_paragraph()
        p.text = t
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(15, 23, 42)

    # Right Bottom Box: ML Lifecycle & Prototype
    life_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(4.5), Inches(5.3), Inches(2.3))
    life_box.fill.solid()
    life_box.fill.fore_color.rgb = RGBColor(240, 253, 244)
    life_box.line.color.rgb = RGBColor(34, 197, 94)
    tf_life = life_box.text_frame
    tf_life.word_wrap = True
    p_l_h = tf_life.paragraphs[0]
    p_l_h.text = "🔄 ML LIFECYCLE & WORKING PROTOTYPE"
    p_l_h.font.bold = True
    p_l_h.font.size = Pt(14)
    p_l_h.font.color.rgb = RGBColor(21, 128, 61)

    p_l1 = tf_life.add_paragraph()
    p_l1.text = "Dataset (10.85M records) ➔ 5-Fold Stratified CV ➔ Holdout Validation ➔ Production Model Deployment"
    p_l1.font.size = Pt(10)
    p_l1.font.bold = True
    p_l1.font.color.rgb = RGBColor(15, 23, 42)

    p_l2 = tf_life.add_paragraph()
    p_l2.text = "✓ Live Working Prototype: Interactive web portal with real-time ML inference (<50ms response time) and automated PDF advice export."
    p_l2.font.size = Pt(10)
    p_l2.font.color.rgb = RGBColor(51, 65, 85)

    # -------------------------------------------------------------
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    setup_slide_header(slide4, "FEASIBILITY AND VIABILITY", 4)

    # 3-Column Infographic Cards
    # Col 1: Technical Feasibility
    f1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(3.9), Inches(4.4))
    f1.fill.solid()
    f1.fill.fore_color.rgb = RGBColor(239, 246, 255)
    f1.line.color.rgb = RGBColor(59, 130, 246)
    tf_f1 = f1.text_frame
    tf_f1.word_wrap = True
    p_f1_h = tf_f1.paragraphs[0]
    p_f1_h.text = "🛠️ TECHNICAL FEASIBILITY"
    p_f1_h.font.bold = True
    p_f1_h.font.size = Pt(13)
    p_f1_h.font.color.rgb = RGBColor(29, 78, 216)

    f1_bullets = [
        "Mature Scikit-learn & Python ecosystem ensures rock-solid stability.",
        "Django framework enables rapid, secure, production-ready web deployment.",
        "Lightweight serialized model artifacts (~400MB) allow sub-50ms inference.",
        "Modular architecture easily handles live soil and meteorological API feeds."
    ]
    for b in f1_bullets:
        p = tf_f1.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(30, 41, 59)

    # Col 2: Operational Feasibility
    f2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.7), Inches(1.4), Inches(3.9), Inches(4.4))
    f2.fill.solid()
    f2.fill.fore_color.rgb = RGBColor(240, 253, 244)
    f2.line.color.rgb = RGBColor(34, 197, 94)
    tf_f2 = f2.text_frame
    tf_f2.word_wrap = True
    p_f2_h = tf_f2.paragraphs[0]
    p_f2_h.text = "🚜 OPERATIONAL FEASIBILITY"
    p_f2_h.font.bold = True
    p_f2_h.font.size = Pt(13)
    p_f2_h.font.color.rgb = RGBColor(21, 128, 61)

    f2_bullets = [
        "Simple, intuitive 2-column input form tailored for rural extension officers.",
        "Farmer receives unambiguous recommendations in local units (kg/ha).",
        "Explainable decision drivers provide clear rationale for every nutrient.",
        "Accessible via standard smartphone and village common service center browsers."
    ]
    for b in f2_bullets:
        p = tf_f2.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(30, 41, 59)

    # Col 3: Challenges & Mitigation
    f3 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(1.4), Inches(3.9), Inches(4.4))
    f3.fill.solid()
    f3.fill.fore_color.rgb = RGBColor(254, 242, 242)
    f3.line.color.rgb = RGBColor(239, 68, 68)
    tf_f3 = f3.text_frame
    tf_f3.word_wrap = True
    p_f3_h = tf_f3.paragraphs[0]
    p_f3_h.text = "🛡️ CHALLENGES & MITIGATIONS"
    p_f3_h.font.bold = True
    p_f3_h.font.size = Pt(13)
    p_f3_h.font.color.rgb = RGBColor(185, 28, 28)

    challs = [
        ("Soil Data Quality Variation", "Automated input range sanitizer & RobustScaler handling."),
        ("Model Reliability Risks", "Strict holdout testing, 5-fold CV & out-of-distribution flags."),
        ("Farmer Trust & Skepticism", "Transparent nutrient explanation instead of black-box advice."),
        ("Regional Agro-Climatic Diversity", "Empirical distribution validation across 32 Indian states.")
    ]
    for c, m in challs:
        p = tf_f3.add_paragraph()
        p.text = f"• {c}: {m}"
        p.font.size = Pt(9.5)
        p.font.color.rgb = RGBColor(30, 41, 59)

    # Bottom Scalability Strip
    scal_box = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.85))
    scal_box.fill.solid()
    scal_box.fill.fore_color.rgb = RGBColor(241, 245, 249)
    scal_box.line.color.rgb = RGBColor(203, 213, 225)
    tf_scal = scal_box.text_frame
    p_scal = tf_scal.paragraphs[0]
    p_scal.text = "📈 FUTURE SCALABILITY ROADMAP:\n[Web Portal Prototype]  ➔  [Multilingual Mobile App]  ➔  [Government Soil Health Card API Sync]  ➔  [IoT Soil Sensor Integration]"
    p_scal.alignment = PP_ALIGN.CENTER
    p_scal.font.bold = True
    p_scal.font.size = Pt(11)
    p_scal.font.color.rgb = RGBColor(15, 23, 42)

    # -------------------------------------------------------------
    # SLIDE 5: IMPACT AND BENEFITS
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    setup_slide_header(slide5, "IMPACT AND BENEFITS", 5)

    # Top Banner Graphic
    imp_banner = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.65))
    imp_banner.fill.solid()
    imp_banner.fill.fore_color.rgb = RGBColor(240, 253, 244)
    imp_banner.line.color.rgb = RGBColor(34, 197, 94)
    tf_ib = imp_banner.text_frame
    p_ib = tf_ib.paragraphs[0]
    p_ib.text = "🌾 DATA-DRIVEN SOIL INPUTS  ➔  🤖 SMART RECOMMENDATION  ➔  📈 OPTIMIZED FARM DECISIONS"
    p_ib.alignment = PP_ALIGN.CENTER
    p_ib.font.bold = True
    p_ib.font.size = Pt(12)
    p_ib.font.color.rgb = RGBColor(21, 128, 61)

    # 4 Impact Cards: FARMERS | ECONOMIC | ENVIRONMENTAL | TECHNOLOGICAL & SOCIAL
    cards_data = [
        ("👨‍🌾 FARMERS", RGBColor(239, 246, 255), RGBColor(59, 130, 246), RGBColor(29, 78, 216), [
            "Eliminates trial-and-error fertilizer selection.",
            "Tailored to specific crop nutritional needs.",
            "Actionable guidance in easy farmer metrics."
        ]),
        ("💰 ECONOMIC", RGBColor(254, 243, 199), RGBColor(245, 158, 11), RGBColor(180, 83, 9), [
            "Reduces expenditure on unnecessary fertilizers.",
            "Maximizes return per rupee spent on inputs.",
            "Prevents crop lodging and nutrient burn."
        ]),
        ("🌍 ENVIRONMENTAL", RGBColor(240, 253, 244), RGBColor(34, 197, 94), RGBColor(21, 128, 61), [
            "Curbs excessive nitrate runoff into groundwater.",
            "Reduces soil salinity and pH degradation.",
            "Promotes sustainable regenerative farming."
        ]),
        ("📱 TECH & SOCIAL", RGBColor(245, 243, 255), RGBColor(139, 92, 246), RGBColor(109, 40, 217), [
            "Democratizes precision agronomy via web access.",
            "Decision-support aid for extension workers.",
            "Scalable API architecture for rural digitization."
        ])
    ]

    for idx, (head, bg_c, border_c, text_c, items) in enumerate(cards_data):
        x = 0.5 + idx * 3.15
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.15), Inches(2.95), Inches(3.7))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_c
        card.line.color.rgb = border_c
        tf_c = card.text_frame
        tf_c.word_wrap = True
        p_h = tf_c.paragraphs[0]
        p_h.text = head
        p_h.font.bold = True
        p_h.font.size = Pt(13)
        p_h.font.color.rgb = text_c

        for item in items:
            p = tf_c.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(30, 41, 59)

    # Bottom Future Potential Strip
    pot_box = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.85))
    pot_box.fill.solid()
    pot_box.fill.fore_color.rgb = RGBColor(241, 245, 249)
    pot_box.line.color.rgb = RGBColor(203, 213, 225)
    tf_pot = pot_box.text_frame
    p_pot = tf_pot.paragraphs[0]
    p_pot.text = "🚀 FUTURE POTENTIAL STRIP:\nRegional Agro-Ecological Models  |  IoT Soil Probes  |  Live Weather Integration  |  Multilingual Voice Interface"
    p_pot.alignment = PP_ALIGN.CENTER
    p_pot.font.bold = True
    p_pot.font.size = Pt(11)
    p_pot.font.color.rgb = RGBColor(15, 23, 42)

    # -------------------------------------------------------------
    # SLIDE 6: RESEARCH AND REFERENCES
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    setup_slide_header(slide6, "RESEARCH AND REFERENCES", 6)

    # Left Column: Research & Data Sources
    r1 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(5.8), Inches(4.4))
    r1.fill.solid()
    r1.fill.fore_color.rgb = RGBColor(248, 250, 252)
    r1.line.color.rgb = RGBColor(59, 130, 246)
    tf_r1 = r1.text_frame
    tf_r1.word_wrap = True
    p_r1_h = tf_r1.paragraphs[0]
    p_r1_h.text = "📚 AGRONOMIC & TECHNICAL REFERENCES"
    p_r1_h.font.bold = True
    p_r1_h.font.size = Pt(13)
    p_r1_h.font.color.rgb = RGBColor(30, 58, 138)

    refs = [
        "• ICAR Guidelines: Soil testing & fertilizer recommendation handbook.",
        "• National Soil Health Card Database: 10.85M survey records analyzed.",
        "• Liebig's Law of the Minimum: Agricultural trace-element limiting factor formulation.",
        "• Scikit-learn & Python Docs: Ensemble methods & RobustScaler pipelines.",
        "• Django Framework Documentation: Scalable RESTful API architecture."
    ]
    for r in refs:
        p = tf_r1.add_paragraph()
        p.text = r
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(30, 41, 59)

    # Right Column: ML Validation Metrics & Scientific Rigor
    r2 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6), Inches(1.4), Inches(6.2), Inches(4.4))
    r2.fill.solid()
    r2.fill.fore_color.rgb = RGBColor(240, 253, 244)
    r2.line.color.rgb = RGBColor(34, 197, 94)
    tf_r2 = r2.text_frame
    tf_r2.word_wrap = True
    p_r2_h = tf_r2.paragraphs[0]
    p_r2_h.text = "📊 RIGOROUS ML VALIDATION METRICS"
    p_r2_h.font.bold = True
    p_r2_h.font.size = Pt(13)
    p_r2_h.font.color.rgb = RGBColor(21, 128, 61)

    metrics = [
        "• Holdout Test Accuracy: 99.75% (Evaluated on 14,945 strictly unseen test samples)",
        "• Macro F1-Score: 0.9961  |  Weighted F1: 0.9975  |  F3-Score (beta=3.0): 0.9960",
        "• Top-2 / Top-3 Accuracy: 100.00%  |  Multi-Class Log Loss: 0.0680",
        "• Regional Distribution Validation: 10,853,209 survey records across 32 States",
        "• Scientific Note: Model learns domain-rule formulations; physical multi-year yield trials remain required for field validation."
    ]
    for m in metrics:
        p = tf_r2.add_paragraph()
        p.text = m
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(15, 23, 42)

    # Bottom QR Code & Project Links
    qr_box = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.85))
    qr_box.fill.solid()
    qr_box.fill.fore_color.rgb = RGBColor(241, 245, 249)
    qr_box.line.color.rgb = RGBColor(203, 213, 225)
    tf_qr = qr_box.text_frame
    p_qr = tf_qr.paragraphs[0]
    p_qr.text = "🔗 PROJECT ARTIFACTS & LINKS:\n[ GitHub: github.com/Shivam-Patel22/SIH-2026-Project- ]   |   [ Live Working Demo: /api/recommendations/generate/ ]   |   [ Full Audit: model_audit.json ]"
    p_qr.alignment = PP_ALIGN.CENTER
    p_qr.font.bold = True
    p_qr.font.size = Pt(11)
    p_qr.font.color.rgb = RGBColor(15, 23, 42)

    output_pptx = 'sih_presentation/SIH_2026_Smart_Fertilizer_Recommendation.pptx'
    prs.save(output_pptx)
    print(f"[+] PPTX saved successfully: {output_pptx}")


def generate_high_res_pdf_presentation():
    """
    Renders high-resolution 16:9 slides using ReportLab and Pillow to produce an export-ready PDF.
    """
    pdf_path = 'sih_presentation/SIH_2026_Smart_Fertilizer_Recommendation.pdf'
    page_width, page_height = 960, 540  # 16:9 widescreen points

    c = canvas.Canvas(pdf_path, pagesize=(page_width, page_height))
    logo_path = 'sih_presentation/extracted_assets/page_1_img_1_Image26.png'
    brain_logo_path = 'sih_presentation/extracted_assets/page_1_img_0_Image7.png'

    # -------------------------------------------------------------
    # SLIDE 1: TITLE PAGE
    # -------------------------------------------------------------
    # Background
    c.setFillColorRGB(1, 1, 1)
    c.rect(0, 0, page_width, page_height, fill=1, stroke=0)

    # Top Right Logo
    if os.path.exists(logo_path):
        c.drawImage(logo_path, page_width - 110, page_height - 90, width=80, height=80, mask='auto')

    # Main Header
    c.setFillColorRGB(15/255, 41/255, 66/255)
    c.setFont("Helvetica-Bold", 26)
    c.drawCentredString(page_width / 2, page_height - 55, "SMART INDIA HACKATHON 2026")

    c.setFillColorRGB(30/255, 58/255, 138/255)
    c.setFont("Helvetica-Bold", 20)
    c.drawCentredString(page_width / 2, page_height - 90, "TITLE PAGE")

    # Project Title Box
    c.setFillColorRGB(240/255, 253/255, 244/255)
    c.setStrokeColorRGB(34/255, 197/255, 94/255)
    c.setLineWidth(1.5)
    c.roundRect(50, page_height - 200, 520, 90, 8, fill=1, stroke=1)

    c.setFillColorRGB(21/255, 128/255, 61/255)
    c.setFont("Helvetica-Bold", 17)
    c.drawString(65, page_height - 145, "AI-Based Smart Fertilizer Recommendation System")

    c.setFillColorRGB(71/255, 85/255, 105/255)
    c.setFont("Helvetica", 11)
    c.drawString(65, page_height - 175, "Precision Nutrient Optimization & Decision-Support Engine for Indian Agriculture")

    # Left Bullets
    c.setFillColorRGB(15/255, 23/255, 42/255)
    c.setFont("Helvetica-Bold", 15)
    c.drawString(50, page_height - 250, "• Problem Statement ID –  [INSERT ACTUAL PS ID]")
    c.drawString(50, page_height - 295, "• Theme –  Agriculture, FoodTech & Rural Development")
    c.drawString(50, page_height - 340, "• Team Name –  [INSERT REGISTERED TEAM NAME]")

    # Right Brain Graphic
    if os.path.exists(brain_logo_path):
        c.drawImage(brain_logo_path, 600, page_height - 390, width=280, height=270, mask='auto')

    # Bottom Strip
    c.setFillColorRGB(241/255, 245/255, 249/255)
    c.setStrokeColorRGB(203/255, 213/255, 225/255)
    c.rect(40, 25, page_width - 80, 42, fill=1, stroke=1)
    c.setFillColorRGB(15/255, 23/255, 42/255)
    c.setFont("Helvetica-Bold", 10)
    c.drawCentredString(page_width / 2, 42, "SOIL CHEMISTRY (N,P,K,pH,Micros)   ➔   ML SOFT-VOTING ENSEMBLE   ➔   PRECISE FERTILIZER DOSAGE   ➔   BETTER CROP YIELD")

    c.showPage()

    # Helper for Slides 2-6
    def draw_slide_template(title_text, template_num):
        c.setFillColorRGB(1, 1, 1)
        c.rect(0, 0, page_width, page_height, fill=1, stroke=0)

        # Top-Left Oval Badge
        c.setStrokeColorRGB(0, 0, 0)
        c.setLineWidth(1.2)
        c.ellipse(40, page_height - 75, 130, page_height - 20)
        c.setFillColorRGB(0, 0, 0)
        c.setFont("Helvetica-Bold", 9)
        c.drawCentredString(85, page_height - 40, "Your")
        c.drawCentredString(85, page_height - 52, "Team")
        c.drawCentredString(85, page_height - 64, "Name")

        # Center Title
        c.setFont("Times-Bold", 24)
        c.drawCentredString(page_width / 2, page_height - 50, title_text)

        # Top-Right Logo
        if os.path.exists(logo_path):
            c.drawImage(logo_path, page_width - 100, page_height - 85, width=70, height=70, mask='auto')

        # Footnote
        c.setFillColorRGB(100/255, 116/255, 139/255)
        c.setFont("Helvetica", 9)
        c.drawString(40, 15, f"@SIH Idea submission- Template {template_num}")

    # -------------------------------------------------------------
    # SLIDE 2: IDEA TITLE / PROPOSED SOLUTION
    # -------------------------------------------------------------
    draw_slide_template("IDEA TITLE", 2)

    c.setFillColorRGB(30/255, 58/255, 138/255)
    c.setFont("Helvetica-Bold", 14)
    c.drawString(40, page_height - 100, "❖ Proposed Solution (Describe your Idea/Solution/Prototype)")

    # 3 Cards: Problem, Solution, Innovation
    card_w = 275
    card_h = 280
    y_card = page_height - 400

    # Problem Card (Red)
    c.setFillColorRGB(254/255, 242/255, 242/255)
    c.setStrokeColorRGB(239/255, 68/255, 68/255)
    c.roundRect(40, y_card, card_w, card_h, 6, fill=1, stroke=1)
    c.setFillColorRGB(185/255, 28/255, 28/255)
    c.setFont("Helvetica-Bold", 12)
    c.drawString(55, y_card + card_h - 25, "⚠️ THE PROBLEM")
    c.setFillColorRGB(51/255, 65/255, 85/255)
    c.setFont("Helvetica", 9.5)
    p_lines = [
        "• Farmers apply fertilizers without knowing",
        "  actual soil nutrient status (guesswork).",
        "• Incorrect choice or dosage inflates input",
        "  costs and reduces farmer profit margins.",
        "• Excessive chemicals cause long-term soil",
        "  acidification & water eutrophication.",
        "• Manual interpretation of 12-parameter",
        "  soil health cards is complex for farmers."
    ]
    for i, line in enumerate(p_lines):
        c.drawString(55, y_card + card_h - 55 - i * 18, line)

    # Solution Card (Blue)
    c.setFillColorRGB(239/255, 246/255, 255/255)
    c.setStrokeColorRGB(59/255, 130/255, 246/255)
    c.roundRect(340, y_card, card_w, card_h, 6, fill=1, stroke=1)
    c.setFillColorRGB(29/255, 78/255, 216/255)
    c.setFont("Helvetica-Bold", 12)
    c.drawString(355, y_card + card_h - 25, "💡 PROPOSED SOLUTION")
    c.setFillColorRGB(51/255, 65/255, 85/255)
    c.setFont("Helvetica", 9.5)
    s_lines = [
        "• AI/ML decision-support system analyzing",
        "  12+ soil parameters (N, P, K, pH, Micros).",
        "• Crop-specific intelligence aligning with",
        "  exact crop stoichiometric requirements.",
        "• Recommends optimum commercial fertilizer",
        "  product + precise application dosage (kg/ha).",
        "• Generates simple, actionable advice with",
        "  transparent soil justification reports."
    ]
    for i, line in enumerate(s_lines):
        c.drawString(355, y_card + card_h - 55 - i * 18, line)

    # Innovation Card (Green)
    c.setFillColorRGB(240/255, 253/255, 244/255)
    c.setStrokeColorRGB(34/255, 197/255, 94/255)
    c.roundRect(640, y_card, card_w, card_h, 6, fill=1, stroke=1)
    c.setFillColorRGB(21/255, 128/255, 61/255)
    c.setFont("Helvetica-Bold", 12)
    c.drawString(655, y_card + card_h - 25, "⭐ INNOVATION & UNIQUENESS")
    c.setFillColorRGB(51/255, 65/255, 85/255)
    c.setFont("Helvetica", 9.5)
    i_lines = [
        "• Integrates Liebig's Law of the Minimum",
        "  across all 6 essential trace micronutrients.",
        "• Soft-Voting Ensemble combining Trees",
        "  (RF, ET), HistGradientBoosting & MLP.",
        "• Validated against empirical distributions",
        "  of 10.85M National Soil Database records.",
        "• Interactive web platform with weather-safe",
        "  application and spray guidance."
    ]
    for i, line in enumerate(i_lines):
        c.drawString(655, y_card + card_h - 55 - i * 18, line)

    # Bottom Pipeline Strip
    c.setFillColorRGB(241/255, 245/255, 249/255)
    c.setStrokeColorRGB(203/255, 213/255, 225/255)
    c.rect(40, 32, page_width - 80, 50, fill=1, stroke=1)
    c.setFillColorRGB(15/255, 23/255, 42/255)
    c.setFont("Helvetica-Bold", 9.5)
    c.drawCentredString(page_width / 2, 60, "SYSTEM PIPELINE: [Soil Data + Crop] ➔ [Data Preprocessing & Scaling] ➔ [Trained ML Model] ➔ [Fertilizer Type + Dosage] ➔ [Farmer Report]")

    c.showPage()

    # -------------------------------------------------------------
    # SLIDE 3: TECHNICAL APPROACH
    # -------------------------------------------------------------
    draw_slide_template("TECHNICAL APPROACH", 3)

    # Left Box: System Architecture
    c.setFillColorRGB(248/255, 250/255, 252/255)
    c.setStrokeColorRGB(59/255, 130/255, 246/255)
    c.roundRect(40, 100, 480, 360, 6, fill=1, stroke=1)

    c.setFillColorRGB(30/255, 58/255, 138/255)
    c.setFont("Helvetica-Bold", 13)
    c.drawString(55, 435, "🏗️ END-TO-END SYSTEM ARCHITECTURE")

    steps = [
        "1. INPUT INGESTION: Soil N, P, K, pH, OC, EC, Zn, B, S, Fe, Mn, Cu + Crop Choice",
        "2. WEB UI LAYER: Responsive Django portal with field-level input sanitization",
        "3. PREPROCESSING: Outlier filtering & RobustScaler fit strictly on training partition",
        "4. FEATURE ENGINEERING: 44 stoichiometric ratios, Liebig factors & soil buffers",
        "5. ML META-ENSEMBLE: Soft-Voting (Random Forest, Extra Trees, HGBoost, Deep MLP)",
        "6. STOICHIOMETRIC ENGINE: Exact stoichiometric balance calculation for kg/ha dosage",
        "7. EXPLAINABLE DASHBOARD: Scientific justification drivers & weather spray warnings"
    ]
    c.setFillColorRGB(30/255, 41/255, 59/255)
    c.setFont("Helvetica", 9)
    for i, st in enumerate(steps):
        c.drawString(55, 400 - i * 36, st)

    # Right Top Box: Technology Stack
    c.setFillColorRGB(239/255, 246/255, 255/255)
    c.setStrokeColorRGB(37/255, 99/255, 235/255)
    c.roundRect(540, 280, 380, 180, 6, fill=1, stroke=1)

    c.setFillColorRGB(29/255, 78/255, 216/255)
    c.setFont("Helvetica-Bold", 12)
    c.drawString(555, 435, "⚙️ TECHNOLOGY STACK")

    tech_lines = [
        "• Frontend: HTML5, Modern Responsive CSS, Vanilla JS",
        "• Backend: Django Framework, Python 3.14 REST API",
        "• Machine Learning: Scikit-learn, NumPy, Pandas, Joblib",
        "• Model Suite: RF (250) + Extra Trees (250) + HGB + MLP",
        "• Database: SQLite (Indexed National Soil DB + Farm Plots)"
    ]
    c.setFillColorRGB(15/255, 23/255, 42/255)
    c.setFont("Helvetica", 9)
    for i, tl in enumerate(tech_lines):
        c.drawString(555, 405 - i * 22, tl)

    # Right Bottom Box: ML Lifecycle & Prototype
    c.setFillColorRGB(240/255, 253/255, 244/255)
    c.setStrokeColorRGB(34/255, 197/255, 94/255)
    c.roundRect(540, 100, 380, 165, 6, fill=1, stroke=1)

    c.setFillColorRGB(21/255, 128/255, 61/255)
    c.setFont("Helvetica-Bold", 12)
    c.drawString(555, 240, "🔄 ML LIFECYCLE & WORKING PROTOTYPE")

    life_lines = [
        "• ML Lifecycle: National Soil Data (10.85M) ➔ 5-Fold Stratified",
        "  Cross-Validation ➔ Final Holdout Test ➔ Live API Deployment",
        "• Working Prototype: Full functional Django web application",
        "  delivering real-time predictions (<50ms) with PDF export."
    ]
    c.setFillColorRGB(15/255, 23/255, 42/255)
    c.setFont("Helvetica", 9)
    for i, ll in enumerate(life_lines):
        c.drawString(555, 210 - i * 20, ll)

    c.showPage()

    # -------------------------------------------------------------
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # -------------------------------------------------------------
    draw_slide_template("FEASIBILITY AND VIABILITY", 4)

    # 3 Columns
    col_w = 275
    col_h = 320
    y_col = 130

    # Col 1: Technical
    c.setFillColorRGB(239/255, 246/255, 255/255)
    c.setStrokeColorRGB(59/255, 130/255, 246/255)
    c.roundRect(40, y_col, col_w, col_h, 6, fill=1, stroke=1)
    c.setFillColorRGB(29/255, 78/255, 216/255)
    c.setFont("Helvetica-Bold", 11.5)
    c.drawString(55, y_col + col_h - 25, "🛠️ TECHNICAL FEASIBILITY")
    c.setFillColorRGB(30/255, 41/255, 59/255)
    c.setFont("Helvetica", 9)
    t_f_lines = [
        "• Mature Scikit-learn & Python stack",
        "  ensures production grade stability.",
        "• Django enables rapid, modular, and",
        "  secure web deployment.",
        "• Lightweight model artifacts (<400MB)",
        "  enable instant real-time inference.",
        "• Easily connects to IoT soil sensors and",
        "  meteorological live APIs."
    ]
    for i, l in enumerate(t_f_lines):
        c.drawString(55, y_col + col_h - 55 - i * 20, l)

    # Col 2: Operational
    c.setFillColorRGB(240/255, 253/255, 244/255)
    c.setStrokeColorRGB(34/255, 197/255, 94/255)
    c.roundRect(340, y_col, col_w, col_h, 6, fill=1, stroke=1)
    c.setFillColorRGB(21/255, 128/255, 61/255)
    c.setFont("Helvetica-Bold", 11.5)
    c.drawString(355, y_col + col_h - 25, "🚜 OPERATIONAL FEASIBILITY")
    c.setFillColorRGB(30/255, 41/255, 59/255)
    c.setFont("Helvetica", 9)
    o_f_lines = [
        "• Clean 2-column input UI designed for",
        "  extension officers and village centers.",
        "• Outputs delivered in standard local",
        "  farmer metrics (kg/ha and commercial bags).",
        "• Explainable justifications build strong",
        "  confidence among farmers.",
        "• Operates on basic mobile browsers with",
        "  low bandwidth requirements."
    ]
    for i, l in enumerate(o_f_lines):
        c.drawString(355, y_col + col_h - 55 - i * 20, l)

    # Col 3: Challenges & Mitigation
    c.setFillColorRGB(254/255, 242/255, 242/255)
    c.setStrokeColorRGB(239/255, 68/255, 68/255)
    c.roundRect(640, y_col, col_w, col_h, 6, fill=1, stroke=1)
    c.setFillColorRGB(185/255, 28/255, 28/255)
    c.setFont("Helvetica-Bold", 11.5)
    c.drawString(655, y_col + col_h - 25, "🛡️ CHALLENGES & MITIGATIONS")
    c.setFillColorRGB(30/255, 41/255, 59/255)
    c.setFont("Helvetica", 8.5)
    cm_lines = [
        "• Soil Data Quality Variance:",
        "  Mitigation: Automated range sanitization &",
        "  RobustScaler handling.",
        "• Model Reliability Risks:",
        "  Mitigation: 5-Fold cross-validation &",
        "  abstention flags for OOD inputs.",
        "• Farmer Trust & Adoption:",
        "  Mitigation: Transparent, explainable rules.",
        "• Regional Soil Diversity:",
        "  Mitigation: Empirical calibration across",
        "  32 Indian states and 735 districts."
    ]
    for i, l in enumerate(cm_lines):
        c.drawString(655, y_col + col_h - 55 - i * 18, l)

    # Bottom Scalability Strip
    c.setFillColorRGB(241/255, 245/255, 249/255)
    c.setStrokeColorRGB(203/255, 213/255, 225/255)
    c.rect(40, 35, page_width - 80, 50, fill=1, stroke=1)
    c.setFillColorRGB(15/255, 23/255, 42/255)
    c.setFont("Helvetica-Bold", 9.5)
    c.drawCentredString(page_width / 2, 60, "FUTURE SCALABILITY: [Web Portal Prototype] ➔ [Multilingual Mobile App] ➔ [Govt SHC Database API] ➔ [IoT Soil Sensor Grid]")

    c.showPage()

    # -------------------------------------------------------------
    # SLIDE 5: IMPACT AND BENEFITS
    # -------------------------------------------------------------
    draw_slide_template("IMPACT AND BENEFITS", 5)

    # Top Graphic
    c.setFillColorRGB(240/255, 253/255, 244/255)
    c.setStrokeColorRGB(34/255, 197/255, 94/255)
    c.rect(40, page_height - 110, page_width - 80, 35, fill=1, stroke=1)
    c.setFillColorRGB(21/255, 128/255, 61/255)
    c.setFont("Helvetica-Bold", 10.5)
    c.drawCentredString(page_width / 2, page_height - 88, "SOIL HEALTH DATA   ➔   INTELLIGENT AI RECOMMENDATION   ➔   SUSTAINABLE FARM OUTCOMES")

    # 4 Impact Cards
    card4_w = 205
    card4_h = 290
    y_card4 = 110

    cards4_data = [
        ("👨‍🌾 FARMERS", 239, 246, 255, 59, 130, 246, 29, 78, 216, [
            "• Eliminates trial-and-error",
            "  fertilizer purchasing.",
            "• Tailored precisely to crop",
            "  nutrient absorption rates.",
            "• Easy-to-follow advice in",
            "  practical farmer units."
        ]),
        ("💰 ECONOMIC", 254, 243, 199, 245, 158, 11, 180, 83, 9, [
            "• Reduces unnecessary input",
            "  costs on surplus fertilizer.",
            "• Enhances crop yield",
            "  efficiency and ROI.",
            "• Prevents crop lodging",
            "  from nitrogen overdose."
        ]),
        ("🌍 ENVIRONMENTAL", 240, 253, 244, 34, 197, 94, 21, 128, 61, [
            "• Curbs nitrate leaching",
            "  into groundwater aquifers.",
            "• Protects soil microbiomes",
            "  and prevents acidification.",
            "• Supports long-term national",
            "  soil conservation."
        ]),
        ("📱 TECH & SOCIAL", 245, 243, 255, 139, 92, 246, 109, 40, 217, [
            "• Democratizes scientific",
            "  precision agronomy.",
            "• Empowers agricultural",
            "  extension field staff.",
            "• Scalable architecture",
            "  ready for rural digital grids."
        ])
    ]

    for idx, (title, r1, g1, b1, r2, g2, b2, r3, g3, b3, items) in enumerate(cards4_data):
        x = 40 + idx * 225
        c.setFillColorRGB(r1/255, g1/255, b1/255)
        c.setStrokeColorRGB(r2/255, g2/255, b2/255)
        c.roundRect(x, y_card4, card4_w, card4_h, 6, fill=1, stroke=1)

        c.setFillColorRGB(r3/255, g3/255, b3/255)
        c.setFont("Helvetica-Bold", 11)
        c.drawString(x + 12, y_card4 + card4_h - 25, title)

        c.setFillColorRGB(30/255, 41/255, 59/255)
        c.setFont("Helvetica", 8.5)
        for j, item in enumerate(items):
            c.drawString(x + 12, y_card4 + card4_h - 55 - j * 20, item)

    # Bottom Future Potential Strip
    c.setFillColorRGB(241/255, 245/255, 249/255)
    c.setStrokeColorRGB(203/255, 213/255, 225/255)
    c.rect(40, 35, page_width - 80, 50, fill=1, stroke=1)
    c.setFillColorRGB(15/255, 23/255, 42/255)
    c.setFont("Helvetica-Bold", 9.5)
    c.drawCentredString(page_width / 2, 60, "FUTURE POTENTIAL: Regional Agro-Ecological Models  |  IoT Soil Probes  |  Live Weather Feeds  |  Multilingual Voice Bot")

    c.showPage()

    # -------------------------------------------------------------
    # SLIDE 6: RESEARCH AND REFERENCES
    # -------------------------------------------------------------
    draw_slide_template("RESEARCH AND REFERENCES", 6)

    # Left Box: References
    c.setFillColorRGB(248/255, 250/255, 252/255)
    c.setStrokeColorRGB(59/255, 130/255, 246/255)
    c.roundRect(40, 100, 420, 360, 6, fill=1, stroke=1)

    c.setFillColorRGB(30/255, 58/255, 138/255)
    c.setFont("Helvetica-Bold", 12)
    c.drawString(55, 435, "📚 AGRONOMIC & TECHNICAL REFERENCES")

    ref_items = [
        "• ICAR Guidelines: Soil Testing & Fertilizer Recommendation",
        "  Handbook (Indian Council of Agricultural Research).",
        "• National Soil Health Card Database: 10,853,209 survey",
        "  records analyzed across 735 districts.",
        "• Liebig's Law of the Minimum: Multi-nutrient limiting",
        "  factor formulation across 6 essential micronutrients.",
        "• Scikit-learn & Python Docs: Voting Classifier, RobustScaler,",
        "  and Stratified K-Fold validation architectures.",
        "• Django Framework Documentation: High-performance",
        "  RESTful web API and model serving."
    ]
    c.setFillColorRGB(30/255, 41/255, 59/255)
    c.setFont("Helvetica", 8.5)
    for i, r in enumerate(ref_items):
        c.drawString(55, 400 - i * 22, r)

    # Right Box: Measured ML Validation Metrics
    c.setFillColorRGB(240/255, 253/255, 244/255)
    c.setStrokeColorRGB(34/255, 197/255, 94/255)
    c.roundRect(480, 100, 440, 360, 6, fill=1, stroke=1)

    c.setFillColorRGB(21/255, 128/255, 61/255)
    c.setFont("Helvetica-Bold", 12)
    c.drawString(495, 435, "📊 MEASURED ML VALIDATION METRICS")

    val_items = [
        "• Holdout Test Accuracy: 99.75% (14,907 / 14,945 strictly unseen test samples)",
        "• Macro F1-Score: 0.9961  |  Weighted F1: 0.9975",
        "• Macro F3-Score (beta=3.0): 0.9960  |  Weighted F3: 0.9975",
        "• Top-2 Accuracy: 100.00%  |  Top-3 Accuracy: 100.00%",
        "• Multi-Class Log Loss: 0.0680  |  Macro Brier Score: 0.0019",
        "• Regional Empirical Validation: 10.85M records across 32 States",
        "• Zero-Data-Leakage Preprocessing: Scaler fitted on Train ONLY",
        "• Scientific Validation Status: PARTIALLY VALIDATED",
        "  (Technical & Synthetic Rules: PASSED | Multi-Year Farm Field Trials: Required)"
    ]
    c.setFillColorRGB(15/255, 23/255, 42/255)
    c.setFont("Helvetica", 8.5)
    for i, v in enumerate(val_items):
        c.drawString(495, 400 - i * 22, v)

    # Bottom QR & Artifact Links
    c.setFillColorRGB(241/255, 245/255, 249/255)
    c.setStrokeColorRGB(203/255, 213/255, 225/255)
    c.rect(40, 35, page_width - 80, 50, fill=1, stroke=1)
    c.setFillColorRGB(15/255, 23/255, 42/255)
    c.setFont("Helvetica-Bold", 9)
    c.drawCentredString(page_width / 2, 60, "PROJECT LINKS: [GitHub: Shivam-Patel22/SIH-2026-Project-]  |  [Working Demo: /api/recommendations/generate/]  |  [Audit: model_audit.json]")

    c.showPage()
    c.save()
    print(f"[+] High-resolution PDF generated: {pdf_path}")


if __name__ == '__main__':
    create_pptx_presentation()
    generate_high_res_pdf_presentation()
